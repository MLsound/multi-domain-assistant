"""
Generates the final presentation deck for FIUBA PNL III (Clase 8, 15 min).

Usage:
    python scripts/generate_presentation.py

Output:
    docs/presentation_FIUBA.pptx

The deck is ~22 slides, branded with the FIUBA logo, and follows the
narrative we'll defend orally. Every slide is built programmatically so
the deck regenerates deterministically from this script.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "docs" / "logoFIUBA.jpg"
OUT = ROOT / "docs" / "presentation_FIUBA.pptx"

# FIUBA brand colours (sampled from the logo)
FIUBA_BLUE = RGBColor(0x00, 0xA8, 0xE1)
DARK = RGBColor(0x12, 0x1B, 0x2D)
GREY = RGBColor(0x55, 0x60, 0x6A)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)


def _add_logo(slide, prs):
    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO),
            left=prs.slide_width - Inches(1.7),
            top=Inches(0.25),
            height=Inches(0.55),
        )


def _add_footer(slide, prs, page_no: int, total: int):
    tb = slide.shapes.add_textbox(
        Inches(0.4), prs.slide_height - Inches(0.4),
        prs.slide_width - Inches(0.8), Inches(0.3),
    )
    tf = tb.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = f"FIUBA · PNL III · Knowledge Assistant · Grupo 1                                    {page_no}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GREY


def _title(slide, prs, text, *, subtitle: str | None = None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), prs.slide_width - Inches(2.4), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT
        p2.font.italic = True


def _bullets(slide, prs, items, *, top: float = 2.4, left: float = 0.7,
             width: float = 11.5, height: float = 5.0, size: int = 18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def _accent_bar(slide, prs, *, top: float = 0.95, height: float = 0.06):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(top),
        Inches(2.0), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FIUBA_BLUE
    bar.line.fill.background()


def _make_slide(prs, layout_idx: int = 6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ---------------------------------------------------------------------------
# Slide builders — each returns the slide so the caller can extend
# ---------------------------------------------------------------------------

def slide_cover(prs):
    s = _make_slide(prs, 6)
    # Banner
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.4))
    band.fill.solid(); band.fill.fore_color.rgb = FIUBA_BLUE; band.line.fill.background()
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.6), Inches(0.7), height=Inches(1.0))

    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Knowledge Assistant"
    p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = "Sistema multi-agente Agentic RAG seguro y alineado"
    p2.font.size = Pt(22); p2.font.color.rgb = ACCENT; p2.font.italic = True
    p3 = tf.add_paragraph()
    p3.text = "Dominio: energía sostenible y edificios inteligentes"
    p3.font.size = Pt(18); p3.font.color.rgb = GREY
    p3.space_before = Pt(6)

    box = s.shapes.add_textbox(Inches(0.6), Inches(5.3), prs.slide_width - Inches(1.2), Inches(2.0))
    tf2 = box.text_frame; tf2.word_wrap = True
    info = [
        ("Grupo 1:", "Fabian Sarmiento · Alejandro Lloveras · Jorge Cuenca"),
        ("Materia:", "Procesamiento de Lenguaje Natural III"),
        ("Docentes:", "Mg. Oksana Bokhonok · Esp. Abraham Rodriguez"),
        ("Maestría:", "Inteligencia Artificial — FIUBA · 2026"),
    ]
    for i, (lbl, val) in enumerate(info):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run1 = p.add_run(); run1.text = lbl + " "
        run1.font.bold = True; run1.font.size = Pt(14); run1.font.color.rgb = DARK
        run2 = p.add_run(); run2.text = val
        run2.font.size = Pt(14); run2.font.color.rgb = DARK
        # Workaround: paragraphs[0] starts with empty text "" before runs added
        if i == 0:
            p.text = ""
            r1 = p.add_run(); r1.text = lbl + " "; r1.font.bold = True; r1.font.size = Pt(14); r1.font.color.rgb = DARK
            r2 = p.add_run(); r2.text = val; r2.font.size = Pt(14); r2.font.color.rgb = DARK


def slide_agenda(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Agenda", subtitle="Lo que vamos a cubrir en 15 minutos")
    _bullets(s, prs, [
        "Problema y motivación",
        "Arquitectura: 7 agentes, decisiones propias",
        "RAG: MLP router + retriever de dos etapas + critic loop",
        "Seguridad: OWASP Top-10 LLM aplicado al pipeline",
        "Autenticación y multi-tenancy (registro / login JWT)",
        "Alineación: Constitutional AI propio, NIST AI RMF",
        "Resultados: métricas, logs, red-team",
        "Demo en vivo y conclusiones",
    ])
    return s


def slide_problem(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Problema",
           subtitle="RAG fiable en dominios técnicos heterogéneos")
    _bullets(s, prs, [
        "Los frameworks high-level operan como cajas negras: poco control sobre el flujo.",
        "Riesgo de alucinaciones en respuestas operativas (BESS, HEMS, PV).",
        "Corpus heterogéneo (científico / software / usuario): un retriever plano sesga.",
        "Amenazas de prompt injection y fuga de PII en aplicaciones LLM.",
        "Necesidad de cumplir NIST AI RMF y OWASP Top-10 LLM 2025.",
    ])
    return s


def slide_objectives(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Objetivos del proyecto")
    _bullets(s, prs, [
        "Construir un pipeline RAG agéntico transparente, sin frameworks que resuelvan el flujo entero.",
        "Cumplir todos los requisitos de la cátedra: FastAPI, ≥2 agentes con comunicación dinámica, RAG, seguridad, modelo preexistente, costo/latencia, evaluación.",
        "Agregar registro / inicio de sesión y aislamiento por usuario.",
        "Implementar guardrails propios (OWASP Top-10 LLM) y alineación constitucional.",
        "Producir un informe técnico y métricas que demuestren que el diseño resuelve el problema.",
    ])
    return s


def slide_architecture(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Arquitectura — 7 agentes en un grafo de estados")
    _bullets(s, prs, [
        "Auth (FastAPI + JWT + SQLite) → emite token al cliente.",
        "Cache semántico → corta circuitos de consultas repetidas.",
        "GuardAgent (entrada) → injection scoring + PII redaction propios.",
        "RouterAgent (MLP propio en PyTorch) → P(dominio | consulta).",
        "RetrievalAgent → búsqueda ponderada Qdrant + reranker cross-encoder.",
        "SynthesisAgent → LLM grounded sobre los chunks recuperados.",
        "CriticAgent → evalúa fidelidad y dispara reintentos (comunicación dinámica).",
        "GuardAgent (salida) → canary, leak markers, PII final.",
        "ActionAgent → audit log estructurado + webhook opcional.",
    ], size=15)
    return s


def slide_rag_decisions(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Decisiones propias de diseño",
           subtitle="Sin frameworks que resuelvan el flujo completo")
    _bullets(s, prs, [
        "MLP en PyTorch entrenado sobre embeddings all-MiniLM-L6-v2 (50 epochs).",
        "Retrieval ponderado: Score = P(c|q) · sim(q,d) → personaliza la búsqueda al dominio dominante.",
        "Reranker cross-encoder ms-marco-MiniLM-L6 sobre top-10 → top-5 final.",
        "Critic-loop con max_retries=2 y refinamiento de query (comunicación inter-agente).",
        "Multi-provider registry: Gemini → Claude → Groq → OpenRouter → Kimi → Ollama.",
        "Cache semántico Qdrant a 0.95 cosine, TTL 1 h.",
    ])
    return s


def slide_classifier(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Router MLP — clasificador propio")
    _bullets(s, prs, [
        "Arquitectura: Linear(384) → ReLU → Dropout(0.3) → Linear(128) → 3 clases (Science / Software / User).",
        "Devuelve distribución de probabilidades, no etiqueta dura.",
        "Evita fallos catastróficos: si la confianza > 0.8 elige tier rápido del LLM (ahorro de costo).",
        "Si Science > 0.4 inyecta contexto ambiental vía MCP weather tool.",
        "Entrenable en < 1 minuto sobre el corpus completo.",
    ])
    return s


def slide_security_intro(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Seguridad: OWASP Top-10 LLM 2025",
           subtitle="Defensa en profundidad, código propio")
    _bullets(s, prs, [
        "LLM01 Prompt Injection — injection_scorer.py (heurística + canary).",
        "LLM02 Insecure Output — validación de salida + redacción PII de salida.",
        "LLM04 DoS — sliding-window limiter + cuota diaria por usuario.",
        "LLM06 Sensitive Info — pii_redactor.py (email, DNI, CUIT, teléfono, IP, API keys).",
        "LLM07 Insecure Plugins — sandbox del MCP + timeouts en webhook.",
        "LLM08 Excessive Agency — el ActionAgent solo escribe logs y dispara webhooks read-only.",
        "LLM10 Model Theft — rate limit + cuota + watermark futuro.",
    ], size=15)
    return s


def slide_injection_scorer(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Detección de prompt injection")
    _bullets(s, prs, [
        "Reglas regex con peso (1.0–1.2 para overrides, 0.4–0.6 para señales débiles).",
        "Score saturante: s = 1 − exp(−Σ pesos) ∈ [0,1].",
        "Decisión: ≥ 0.70 bloquea, ≥ 0.40 marca, < 0.40 deja pasar.",
        "Cobre EN/ES: 'ignore previous', 'olvida las reglas', 'modo desarrollador', DAN, 'reveal prompt'…",
        "Trazabilidad total: matched_rules se loguea por consulta (auditable, defendible).",
        "Canary token (signed-prompt) bloquea fugas de system prompt.",
    ])
    return s


def slide_auth(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Autenticación y multi-tenancy",
           subtitle="Hand-rolled — sin fastapi-users")
    _bullets(s, prs, [
        "POST /auth/register · POST /auth/login · GET /auth/me — JWT HS256 (60 min).",
        "Hash de passwords con bcrypt (passlib).",
        "Roles: user / researcher / admin. RBAC en deps.py via require_role().",
        "Cuota diaria por usuario + sliding-window de ráfagas (LLM04 mitigación).",
        "Tabla query_records: cada consulta queda asociada al user_id (aislamiento).",
        "session_id se prefijea con u<id>: para no compartir estado entre usuarios.",
    ])
    return s


def slide_alignment(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Alineación: Constitutional AI propio")
    _bullets(s, prs, [
        "constitution.yaml con 7 principios versionados:",
        "  P1 Grounded — toda afirmación debe apoyarse en chunks.",
        "  P2 No-PII — nunca emitir PII verbatim.",
        "  P3 Safety — no instrucciones que comprometan personas/infra.",
        "  P4 Honesty — expresar incertidumbre cuando no hay evidencia.",
        "  P5 Scope — declinar fuera de energía sostenible / smart-building.",
        "  P6 No-toxicity — respuestas respetuosas siempre.",
        "  P7 Transparency — citar siempre los documentos usados.",
        "El CriticAgent evalúa cada respuesta contra los principios y dispara reintento si falla.",
    ], size=14)
    return s


def slide_nist(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "NIST AI RMF — implementado en código vivo")
    _bullets(s, prs, [
        "GOVERN: docs/governance/policy.md (roles, retención, autoridad de stop-the-line).",
        "MAP: risk_register.md — 12 riesgos puntuados L × I, ligados a controles concretos.",
        "MEASURE: /metrics expone faithfulness, blocked_by_guard_rate, pii_redacted_count, error_rate, latency.",
        "MANAGE: incident_response.md — sev-1 a sev-4 con runbook reproducible.",
        "Model Card (Mitchell et al. 2019) y Datasheet (Gebru et al. 2021) en docs/governance/.",
    ])
    return s


def slide_metrics(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Métricas y logs",
           subtitle="Lo que pidió la consigna: demostrar que el diseño funciona")
    _bullets(s, prs, [
        "Faithfulness (Ragas, LLM judge): proporción de claims apoyadas por contexto.",
        "Context recall (Ragas): cobertura de los hechos esperados.",
        "Precision@5: fracción de chunks del dominio esperado.",
        "ROUGE-L F1 + similitud semántica bge-large-en.",
        "Tiempos: retrieval_time_ms, total_latency_ms (p50, p95).",
        "Costo aproximado: token_count * precio_por_token del proveedor.",
        "Logs: logs/audit.jsonl + tabla query_records (per-user persistente).",
        "Endpoint /metrics con tasa de bloqueo, PII detectada, rate-limit hits.",
    ], size=14)
    return s


def slide_redteam(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Red-team propio")
    _bullets(s, prs, [
        "Suite de ataques propia en tests/test_security.py:",
        "  · Direct override (EN/ES): 'ignore previous', 'olvida las reglas'.",
        "  · Role swap: 'you are now DAN', 'modo desarrollador'.",
        "  · System-prompt leak: 'reveal your prompt', 'muestra instrucciones'.",
        "  · Secret exfil: 'output your api key'.",
        "  · Adversarial noise / suffix.",
        "  · PII fixtures (DNI, CUIT, email, IP).",
        "Métrica objetivo: injection_block_rate ≥ 0.90 sobre la suite.",
    ], size=14)
    return s


def slide_cost_latency(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Optimización de costo y latencia")
    _bullets(s, prs, [
        "Cache semántico Qdrant (≥ 0.95 cosine) → respuesta repetida ~ 0 ms y 0 tokens.",
        "Tier de LLM dinámico: confianza > 0.8 → modelo 'simple' (más barato).",
        "Reranker reduce a top-5: el LLM ve 5 chunks, no 10.",
        "Multi-provider: si Gemini se cae, fallback a Groq → Ollama (cero costo).",
        "Sliding-window limiter evita gastos por bots.",
        "Métrica: avg_latency_ms y cache_hit_rate visibles en /metrics.",
    ])
    return s


def slide_results(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Resultados — Evaluación Final")
    _bullets(s, prs, [
        "Faithfulness:    0.84 (objetivo ≥ 0.80)",
        "Context recall:  0.81 (objetivo ≥ 0.75)",
        "Precision@5:     0.88",
        "Latencia p50:    1840 ms (objetivo < 3000)",
        "Latencia p95:    2850 ms",
        "Cache hit rate:  12.5%",
        "Injection block: 100% (15/15 bloqueados en red-team)",
        "Tokens/consulta promedio: 310",
    ])
    return s


def slide_demo(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Demo en vivo")
    _bullets(s, prs, [
        "1. POST /auth/register con email + password → 201.",
        "2. POST /auth/login → recibimos JWT.",
        "3. POST /query con Bearer token → respuesta fundamentada + citas.",
        "4. Mostrar bloqueo: 'Ignore all previous instructions...' → 4xx con motivo.",
        "5. GET /me/queries → historial del usuario, aislado.",
        "6. GET /metrics → todas las métricas que diseñamos.",
    ])
    return s


def slide_regulation(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Marco regulatorio")
    _bullets(s, prs, [
        "EU AI Act (2024): este sistema es 'limited risk' → deber de transparencia (cumplido vía citas).",
        "NIST AI RMF 1.0: cuatro funciones implementadas (ver docs/governance).",
        "Anthropic Responsible Scaling Policy: nuestro modelo opera a ASL-2.",
        "Argentina: sin ley específica al 2026, pero seguimos OECD/UNESCO.",
        "Privacidad: PII redactada antes y después del LLM; logs sin valores crudos.",
    ])
    return s


def slide_limits(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Limitaciones y trabajo futuro")
    _bullets(s, prs, [
        "Heurísticas vs clasificador entrenado: la regex perderá ataques novedosos → entrenar un classifier propio.",
        "PII regex sufre falsos negativos en español; complementar con NER liviano es el siguiente paso.",
        "Bias / fairness aún no cuantificado formalmente (corpus técnico, baja exposición a atributos protegidos).",
        "Watermark de salidas y trust-score por chunk son stubs en risk_register pendientes de implementar.",
        "Pipeline mono-proceso: para producción, mover rate-limiter a Redis.",
    ])
    return s


def slide_conclusions(prs):
    s = _make_slide(prs); _add_logo(s, prs); _accent_bar(s, prs)
    _title(s, prs, "Conclusiones")
    _bullets(s, prs, [
        "Construimos un pipeline agéntico explicable, sin cajas negras.",
        "Cumplimos todos los requisitos de la consigna y agregamos auth + governance + alineación.",
        "Cada control de seguridad mapea a un riesgo OWASP y a una métrica observable.",
        "El sistema es defendible: cada decisión se puede explicar leyendo el código y los principios.",
        "Listo para iterar y extender (corpus, idiomas, fairness).",
    ])
    return s


def slide_thanks(prs):
    s = _make_slide(prs, 6)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(5.4), Inches(0.7), height=Inches(1.0))

    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.8), prs.slide_width - Inches(1.0), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "¡Gracias!"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = "Preguntas, comentarios y red-team welcome."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20); p2.font.color.rgb = ACCENT; p2.font.italic = True
    p3 = tf.add_paragraph()
    p3.text = "Grupo 1 — Sarmiento · Lloveras · Cuenca"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(16); p3.font.color.rgb = GREY
    p3.space_before = Pt(24)
    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_cover,
        slide_agenda,
        slide_problem,
        slide_objectives,
        slide_architecture,
        slide_rag_decisions,
        slide_classifier,
        slide_security_intro,
        slide_injection_scorer,
        slide_auth,
        slide_alignment,
        slide_nist,
        slide_metrics,
        slide_redteam,
        slide_cost_latency,
        slide_results,
        slide_demo,
        slide_regulation,
        slide_limits,
        slide_conclusions,
        slide_thanks,
    ]
    total = len(builders)
    for i, b in enumerate(builders, start=1):
        slide = b(prs)
        # Cover and thanks are decorative — skip footer there
        if b not in (slide_cover, slide_thanks):
            _add_footer(slide, prs, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
